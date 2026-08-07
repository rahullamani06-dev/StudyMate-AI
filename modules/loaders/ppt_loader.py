"""
ppt_loader.py

Purpose:
--------
Extract text from Microsoft PowerPoint (.pptx) files.

Author: Rahul I Lamani
Project: StudyMate AI
"""

from pathlib import Path
from pptx import Presentation


class PPTLoader:
    """
    A class to load and extract text from PPTX files.
    """

    def __init__(self, file_path: str):
        self.file_path = Path(file_path)

    def extract_text(self) -> str:
        """
        Extract text from all slides.

        Returns:
            str: Complete extracted text.
        """

        if not self.file_path.exists():
            raise FileNotFoundError(f"File not found: {self.file_path}")

        extracted_text = []

        try:
            presentation = Presentation(self.file_path)

            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()

                        if text:
                            extracted_text.append(text)

            return "\n".join(extracted_text)

        except Exception as error:
            raise Exception(f"Error reading PPTX file: {error}")